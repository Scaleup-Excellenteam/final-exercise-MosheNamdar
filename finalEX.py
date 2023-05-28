import os
from pptx import Presentation
import openai
import asyncio
import json
from os import path
from dotenv import load_dotenv
from datetime import datetime, timedelta

# Load environment variables from .env file
load_dotenv()

# Retrieve OpenAI API key from environment variable
KEY = os.environ.get("KEY")

# Variable to keep track of the number of API requests made
numbers_of_requests = 0


def save_answer_on_json_file(answer, json_file_name):
    """
    Save the answer in a JSON file.

    Args:
        answer (str): The answer to be saved.
        json_file_name (str): The name of the JSON file.
    """
    json.dump(answer, open(json_file_name, "w"))


async def generate_answer_gpt(slides):
    """
    Generates an answer using OpenAI's GPT-3.5-turbo model.

    Args:
        slides (list): A list of slide texts.

    Returns:
        str: The generated answer.
    """
    openai.api_key = KEY
    global numbers_of_requests

    # Build the request message
    request = "Hi, could you please summarize the following slides for me?\n"
    answer = ""
    for i, slide in enumerate(slides):
        request += slide

        try:
            # Make an API call to generate a response
            response = await asyncio.to_thread(openai.ChatCompletion.create,
                                               model="gpt-3.5-turbo",
                                               messages=[
                                                   {"role": "system", "content": request},
                                               ])

            # Extract the answer from the API response
            slide_answer = response.choices[0].message.content.strip()
            answer += slide_answer

        except Exception as e:
            # Handle errors raised while processing a slide
            slide_error_message = f"Error processing slide {i + 1}: {str(e)}\n\n"
            answer += slide_error_message

        numbers_of_requests += 1
        prev_request_time = datetime.now()

        # Pause after every third request to stay within rate limits
        if i < len(slides) - 1 and numbers_of_requests % 3 == 0:
            time_left = timedelta(minutes=1) - (datetime.now() - prev_request_time)

            if time_left.total_seconds() > 0:
                await asyncio.sleep(time_left.total_seconds())

    return answer


def read_presentation(pptx_path):
    """
    Reads a PowerPoint presentation and extracts the text from each slide.

    Args:
        pptx_path (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of slide texts.
    """
    file = Presentation(pptx_path)
    text = []
    for slide in file.slides:
        text_slide = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for ph in shape.text_frame.paragraphs:
                    text_slide += ph.text
        text.append(text_slide)

    # print(text)
    return text


def get_output_file_name(pptx_file_name):
    """
    Generates the output file name based on the PowerPoint file name.

    Args:
        pptx_file_name (str): The name of the PowerPoint file.

    Returns:
        str: The output file name.
    """
    output_file_name = path.splitext(pptx_file_name)[0] + ".json"
    return output_file_name


async def main():
    """
    Main function to execute the program.
    """
    path_of_pptx = input("please enter a path: ")

    presentation_text = read_presentation(path_of_pptx)

    answer = await generate_answer_gpt(presentation_text)
    # print(answer)

    json_file_name = get_output_file_name(path.basename(path_of_pptx))

    save_answer_on_json_file(answer, json_file_name)


if __name__ == '__main__':
    asyncio.run(main())
