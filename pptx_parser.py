from pptx import Presentation


def read_presentation(pptx_path):
    """
    Reads a PowerPoint presentation and extracts the text from each slide.

    Args:
        pptx_path (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of slide texts.
    """
    file = Presentation(pptx_path)
    text = []
    for slide in file.slides:
        text_slide = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for ph in shape.text_frame.paragraphs:
                    text_slide += ph.text
        text.append(text_slide)

    # print(text)
    return text
